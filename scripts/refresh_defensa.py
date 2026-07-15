#!/usr/bin/env python3
"""Refresco CONSERVADOR (estilo UNIR) de la presentación de defensa (sobre COPIA).

Filosofía: institucional y sobrio — NO se cambia paleta, fondo ni tipografía.
Solo se añade:
- Portada: una regla de acento fina (navy) bajo el título.
- Todas las slides: pie de página discreto en gris claro con
  "TFM · UNIR · Junio 2026" (izq.) y numeración "N / total" (der.).

NO toca el original: lee SRC y escribe DST.
Uso:  ../.venv_thesis/bin/python scripts/refresh_defensa.py
"""

from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
# SRC = original pre-estilo (fuente); DST = archivo principal (con el refresco).
SRC = REPO / "docs/defensa/Presentacion_Defensa_TFM_base.pptx"
DST = REPO / "docs/defensa/Presentacion_Defensa_TFM.pptx"

NAVY = RGBColor(0x1F, 0x4E, 0x79)  # acento institucional (ya usado en el deck)
GREY = RGBColor(0x8A, 0x8A, 0x8A)  # pie de página discreto
FOOTER_LEFT = "TFM · UNIR · Junio 2026"
MARK = "footer_unir"  # marca de idempotencia


def _has_marker(slide) -> bool:
    for sh in slide.shapes:
        if sh.name == MARK:
            return True
    return False


def add_footer(slide, page: int, total: int) -> None:
    # izquierda
    left = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(6.0), Inches(0.35))
    left.name = MARK
    p = left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = FOOTER_LEFT
    r.font.size = Pt(9)
    r.font.color.rgb = GREY
    # derecha (numeración)
    right = slide.shapes.add_textbox(Inches(11.0), Inches(7.05), Inches(1.9), Inches(0.35))
    right.name = MARK
    pr = right.text_frame.paragraphs[0]
    pr.alignment = PP_ALIGN.RIGHT
    rr = pr.add_run()
    rr.text = f"{page} / {total}"
    rr.font.size = Pt(9)
    rr.font.color.rgb = GREY


def add_title_accent_rule(slide) -> None:
    """Banda de acento institucional (navy) en los bordes de la portada.

    Estilo membrete: una banda fina arriba y otra abajo. No colisiona con el
    bloque de título centrado.
    """
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(13.333), Pt(9)
    )
    bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(7.4), Inches(13.333), Pt(9)
    )
    for band in (top, bottom):
        band.name = MARK
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()


def main() -> int:
    shutil.copy(SRC, DST)
    prs = Presentation(DST)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if _has_marker(slide):
            continue
        # numeración en todas MENOS la portada (convención académica)
        if i > 1:
            add_footer(slide, i, total)
        else:
            add_title_accent_rule(slide)
    prs.save(DST)
    print(f"refresco UNIR aplicado -> {DST.relative_to(REPO)} ({total} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
