#!/usr/bin/env python3
"""Embebe videos de simulación en los *_base.pptx (fuente reproducible).

Motivación: poder PROYECTAR las simulaciones desde la propia presentación.
El reel de lenguaje natural (la contribución más nueva, Entrega 4) no estaba
en ninguna presentación; este script lo embebe en el slide correspondiente de
ambas, reemplazando la imagen estática por el video (con poster frame).

Diseño:
- Opera sobre los *_base.pptx (las canónicas se regeneran luego con los
  scripts de estilo, que preservan el video al copiar el fichero).
- Idempotente: si el slide ya tiene un shape de video con la marca, no hace nada.
- El poster frame se extrae con ffmpeg en un instante con contenido visible.

Uso:  ../.venv_thesis/bin/python scripts/embed_sim_videos.py
Requiere: ffmpeg en PATH.
"""
from __future__ import annotations

import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

REPO = Path(__file__).resolve().parents[1]
RESULTS = REPO / "experiments/results"
MARK = "sim_video"  # marca de idempotencia (prefijo del nombre del shape)

# (deck_base, slide_idx 1-based, video, imagen_a_reemplazar, poster_t_segundos)
JOBS = [
    # Se usa el corte de ~39s (5 instrucciones a 4.5x) para no comerse el
    # tiempo de la defensa; el reel completo queda en videos_proyeccion/.
    ("docs/entrega3/Presentacion_Defensa_TFM_base.pptx", 13,
     "language_reel/language_reel_corto.mp4", "Picture 4", 3.0),
    ("docs/entrega3/Presentacion_Robotica_IA_base.pptx", 13,
     "language_reel/language_reel_corto.mp4", "Picture 3", 3.0),
]


def _poster(video: Path, t: float, out: Path) -> Path | None:
    """Extrae un frame en t segundos como poster (PNG). None si falla."""
    try:
        subprocess.run(
            ["ffmpeg", "-y", "-ss", str(t), "-i", str(video),
             "-frames:v", "1", "-q:v", "2", str(out)],
            check=True, capture_output=True)
        return out if out.exists() else None
    except Exception:
        return None


def _find_shape(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _has_video_mark(slide) -> bool:
    return any(sh.name.startswith(MARK) for sh in slide.shapes)


def embed_one(deck: Path, idx: int, video: Path, repl_name: str,
              poster_t: float, tmp: Path) -> str:
    prs = Presentation(deck)
    slide = list(prs.slides)[idx - 1]
    if _has_video_mark(slide):
        return f"  [skip] {deck.name} s{idx}: ya tiene video embebido"

    # footprint: tomar el de la imagen a reemplazar; si no, región derecha 16:9
    repl = _find_shape(slide, repl_name)
    if repl is not None:
        left, top, w, h = repl.left, repl.top, repl.width, repl.height
        # eliminar la imagen estática (el poster del video la sustituye)
        repl._element.getparent().remove(repl._element)
    else:
        left, top, w, h = Inches(7.0), Inches(1.8), Inches(5.8), Inches(3.26)

    # encajar 16:9 dentro del footprint disponible, centrado
    vw, vh = 16, 9
    fit_w = w
    fit_h = int(fit_w * vh / vw)
    if fit_h > h:
        fit_h = h
        fit_w = int(fit_h * vw / vh)
    left = left + (w - fit_w) // 2
    top = top + (h - fit_h) // 2

    poster = _poster(video, poster_t, tmp / f"poster_{idx}.png")
    movie = slide.shapes.add_movie(
        str(video), left, top, fit_w, fit_h,
        poster_frame_image=str(poster) if poster else None,
        mime_type="video/mp4")
    movie.name = f"{MARK}_{video.stem}"
    prs.save(deck)
    return (f"  [ok]   {deck.name} s{idx}: + {video.name} "
            f"({Emu(fit_w).inches:.1f}x{Emu(fit_h).inches:.1f}in"
            f"{', poster' if poster else ', sin poster'})")


def main() -> int:
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        for deck_rel, idx, vid_rel, repl, t in JOBS:
            deck = REPO / deck_rel
            video = RESULTS / vid_rel
            if not video.exists():
                print(f"  [WARN] no existe el video: {video}")
                continue
            print(embed_one(deck, idx, video, repl, t, tmp))
    print("\nListo. Ahora regenera las canónicas con los scripts de estilo:")
    print("  ../.venv_thesis/bin/python scripts/restyle_divulga_pastel.py")
    print("  ../.venv_thesis/bin/python scripts/refresh_defensa.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
