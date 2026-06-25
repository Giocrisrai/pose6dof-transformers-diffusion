#!/usr/bin/env python3
"""Añade botones CLICABLES para abrir las demos en la divulgativa (_base).

- Slide 11 'DEMO 1 — Ustedes mandan': botón → http://127.0.0.1:7860 (Gradio en
  vivo; requiere `python scripts/demo_charla.py` corriendo).
- Slide 13 'DEMO 2 — En vivo en el simulador': botón → vídeo E2E grabado
  (ruta relativa al .pptx; plan B si la demo en vivo no está disponible).

Los botones se crean en el _base con colores que sobreviven al restyle (fill
de la paleta Vibrante tech, texto blanco). Idempotente por nombre de shape.
Correr ANTES de restyle_robotica_ia.py.

Uso:  ../.venv_thesis/bin/python scripts/add_demo_buttons.py
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

REPO = Path(__file__).resolve().parents[1]
DECK = REPO / "docs/entrega3/Presentacion_Robotica_IA_base.pptx"

CYAN = RGBColor(0x22, 0xD3, 0xEE)
LIME = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MARK = "demo_btn"

# (idx_slide_1based, marcador_titulo, texto_boton, url/href, color)
BUTTONS = [
    (11, "DEMO 1", "▶  Abrir demo interactiva  ·  localhost:7860",
     "http://127.0.0.1:7860", CYAN),
    (13, "DEMO 2", "▶  Abrir demo E2E grabada",
     "videos_proyeccion/03_e2e_con_telemetria.mp4", LIME),
]


def _slide_by_title(prs, mark):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and mark in sh.text_frame.text:
                return s
    return None


def _has(slide):
    return any(sh.name == MARK for sh in slide.shapes)


def add_button(slide, text, href, color):
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.7), Inches(6.45), Inches(6.4), Inches(0.7))
    btn.name = MARK
    btn.fill.solid(); btn.fill.fore_color.rgb = color
    btn.line.fill.background()
    btn.click_action.hyperlink.address = href     # toda la forma es clicable
    tf = btn.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE


def main() -> int:
    prs = Presentation(DECK)
    done = 0
    for idx, mark, text, href, color in BUTTONS:
        slide = _slide_by_title(prs, mark)
        if slide is None:
            print(f"  [warn] no encontré slide con '{mark}'")
            continue
        if _has(slide):
            print(f"  [skip] '{mark}' ya tiene botón")
            continue
        add_button(slide, text, href, color)
        print(f"  [ok] botón en '{mark}' -> {href}")
        done += 1
    if done:
        prs.save(DECK)
    print("Ahora regenera la canónica: scripts/restyle_robotica_ia.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
