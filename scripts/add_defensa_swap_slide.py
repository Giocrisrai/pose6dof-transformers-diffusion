#!/usr/bin/env python3
"""Inserta en la defensa (_base) un slide sobre percepción intercambiable /
comercializable, tras 'Honestidad y limitaciones'. Estilo UNIR del deck.

Convierte la limitación de licencia (FoundationPose NC) en oportunidad: el
bloque de percepción es intercambiable por uno open-license sin tocar el resto.

Idempotente (no duplica si ya existe). Correr; luego refresh_defensa.py.
Uso:  ../.venv_thesis/bin/python scripts/add_defensa_swap_slide.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
DECK = REPO / "docs/entrega3/Presentacion_Defensa_TFM_base.pptx"

NAVY = RGBColor(0x1F, 0x4E, 0x79)
INK = RGBColor(0x22, 0x22, 0x22)
TITLE = "Percepción intercambiable → comercializable"
AFTER_TITLE = "Honestidad y limitaciones"

BULLETS = [
    "La percepción está desacoplada: protocolo PoseEstimator (predict_pose → pose en SE(3)) con is_commercializable(); el resto del pipeline no depende del estimador",
    "FoundationPose es de uso no comercial → se sustituye por un bloque open-license cambiando solo pose_method, sin tocar nada más",
    "FreeZeV2.1 (Apache-2.0) ganó el BOP Challenge 2024: estado del arte y comercializable. Evidencia propia: FreeZeV2 a solo −3 pp AUC (exploración #3)",
    "O GDR-Net++ (Apache-2.0) fine-tuneado al catálogo de la planta → cierra la brecha en un dominio acotado",
    "El aporte del TFM (puente matemático SE(3) + difusión + reproducibilidad) se conserva: es independiente del estimador concreto",
]


def _index_of(prs, title_substr):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if sh.has_text_frame and title_substr in sh.text_frame.text:
                return i
    return None


def main() -> int:
    prs = Presentation(DECK)
    if _index_of(prs, "Percepción intercambiable") is not None:
        print("  [skip] el slide ya existe")
        return 0
    after = _index_of(prs, AFTER_TITLE)
    if after is None:
        raise SystemExit(f"No encontré '{AFTER_TITLE}'")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Título
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1.0))
    tb.text_frame.word_wrap = True
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = TITLE
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    # Divisor
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35),
                                 Inches(12.33), Pt(2.5))
    div.fill.solid(); div.fill.fore_color.rgb = NAVY; div.line.fill.background()
    # Bullets
    bb = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(12.33), Inches(5.0))
    btf = bb.text_frame; btf.word_wrap = True
    for i, b in enumerate(BULLETS):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run(); run.text = "•  " + b
        run.font.size = Pt(17); run.font.color.rgb = INK

    # mover el slide recién creado a la posición after+1
    lst = prs.slides._sldIdLst
    new = list(lst)[-1]
    lst.remove(new)
    lst.insert(after + 1, new)
    prs.save(DECK)
    print(f"  [ok] insertado '{TITLE}' como slide {after + 2}")
    print("Ahora regenera la canónica: scripts/refresh_defensa.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
