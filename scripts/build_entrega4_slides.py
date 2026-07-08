#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""build_entrega4_slides.py — SP-3 (Entrega 4)

Inserta UNA diapositiva nueva sobre la capacidad de "bin picking guiado por
lenguaje natural" en cada uno de los dos decks (defensa y divulgativa),
replicando el estilo propio de cada deck y de forma idempotente:

  - Si ya existe una slide con el título-marcador, NO se duplica: se imprime un
    mensaje y se omite ese deck.
  - La slide se construye sobre layout 6 ("Blank") y se reordena a la posición
    objetivo (índice 12 -> queda como slide 13 en ambos decks).

Requiere python-pptx (disponible en ../.venv_thesis). Ejecutar como:
    ../.venv_thesis/bin/python scripts/build_entrega4_slides.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# --- Rutas (relativas a la raíz del repo, que es el cwd esperado) -----------
DECK_DEFENSA = "docs/entrega3/Presentacion_Defensa_TFM.pptx"
DECK_DIVULGA = "docs/entrega3/Presentacion_Robotica_IA.pptx"
ASSET_IMG = "docs/entrega4/assets/fig_language_pick.png"

# --- Marcadores de idempotencia ---------------------------------------------
MARKER_DEFENSA = "Bin picking guiado por lenguaje natural"
MARKER_DIVULGA = "Le hablas"

TARGET_INDEX = 12  # 0-based; queda como slide 13 (después de la slide 12 1-based)

# --- Notas del ponente (verbatim) -------------------------------------------
NOTE_DEFENSA = (
    "Como capa de Entrega 4, el pipeline ahora entiende lenguaje natural. "
    "Ejemplo: «dame el cubo rojo de la izquierda» → el sistema ancla la "
    "instrucción a un objeto por sus atributos (color, forma, tamaño) y su "
    "relación espacial, y lo agarra con la misma cadena de percepción y "
    "planificación. Por defecto usa un parser determinista en español e "
    "inglés, reproducible y sin red; opcionalmente se enchufa un modelo de "
    "lenguaje local. Validación honesta: la selección del objeto correcto "
    "acierta el 100 % en banco controlado (90 instrucciones puras + 9 en "
    "simulación) y la robustez ante distractores está en los experimentos 16 "
    "a 26; el agarre es cinemático, validado por proximidad pinza-objeto de 4 "
    "mm y convergencia de la cinemática inversa. Es opt-in: extiende el "
    "pipeline base, no lo altera."
)
NOTE_DIVULGA = (
    "Momento «háblale al robot»: le pido en lenguaje natural el objeto y lo "
    "coge. Recalcar que es el mismo sistema de antes —ojos, cerebro y brazo— "
    "ahora con comprensión de lenguaje. Tono cercano; no entrar en métricas "
    "aquí, eso es para la defensa técnica."
)


# ---------------------------------------------------------------------------
def _slide_titles(prs):
    """Lista de títulos (texto de la primera caja con texto) por slide."""
    titles = []
    for s in prs.slides:
        title = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                title = sh.text_frame.text.strip()
                break
        titles.append(title)
    return titles


def _has_marker(prs, marker):
    """True si alguna slide contiene el marcador en cualquier caja de texto."""
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text:
                return True
    return False


_NOTES_BODY_SP_XML = (
    '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<p:nvSpPr>'
    '<p:cNvPr id="2" name="Notes Placeholder 1"/>'
    '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
    '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr>'
    '</p:nvSpPr>'
    '<p:spPr/>'
    '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
    '</p:sp>'
)


def _ensure_notes_text_frame(notes_slide):
    """Devuelve el notes_text_frame, inyectando un placeholder BODY en el XML
    de la notes slide si el deck carece de él en su notes master (en cuyo caso
    python-pptx devuelve None).
    """
    tf = notes_slide.notes_text_frame
    if tf is not None:
        return tf
    from pptx.oxml import parse_xml
    sp = parse_xml(_NOTES_BODY_SP_XML)
    spTree = notes_slide.shapes._spTree
    spTree.append(sp)
    tf = notes_slide.notes_text_frame
    if tf is None:
        raise RuntimeError("No se pudo crear el text frame de notas")
    return tf


def set_notes(prs, marker_substring, note_text):
    """Localiza la primera slide cuyo texto contiene marker_substring y le
    asigna note_text como notas del ponente (idempotente: sobrescribe).

    Lanza ValueError si no encuentra ninguna slide con el marcador.
    """
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and marker_substring in sh.text_frame.text:
                notes = s.notes_slide  # crea la notes slide si no existe
                _ensure_notes_text_frame(notes).text = note_text
                return s
    raise ValueError(
        f"No se encontró ninguna slide con el marcador {marker_substring!r}"
    )


def _move_last_slide_to(prs, target_index):
    """Mueve la última slide (recién añadida) a target_index (0-based)."""
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    new = sld_ids[-1]
    sldIdLst.remove(new)
    sldIdLst.insert(target_index, new)


# ---------------------------------------------------------------------------
def build_defensa_slide(prs):
    """Construye la slide de lenguaje natural en el estilo del deck de defensa.

    Estilo del deck (slides de contenido 12-14): caja-título + rectángulo
    divisor azul + caja de bullets con prefijo "•  ".
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    azul = RGBColor(0x1F, 0x4E, 0x79)

    # Título
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = MARKER_DEFENSA
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = azul

    # Divisor
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.4), Inches(12.3), Pt(3)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = azul
    div.line.fill.background()

    # Bullets
    bullets = [
        "«dame el cubo rojo de la izquierda» → el robot selecciona y agarra el objeto descrito",
        "Parser determinista ES/EN + modelo de lenguaje local enchufable (con fallback)",
        "Anclaje por atributos (color/forma/tamaño) y relación espacial",
        "Open-license, corre en portátil; opt-in sobre el pipeline base (no lo altera)",
        "Validado E2E en CoppeliaSim: selección 100 % (pura n=90 / sim n=9), agarre 4 mm, IK convergente",
    ]
    bb = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(6.3), Inches(4.8))
    btf = bb.text_frame
    btf.word_wrap = True
    for i, b in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        run = para.add_run()
        run.text = "•  " + b
        run.font.size = Pt(17)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    # Imagen
    slide.shapes.add_picture(ASSET_IMG, Inches(7.1), Inches(1.8), width=Inches(5.8))

    return slide


def build_divulga_slide(prs):
    """Construye la slide de lenguaje natural en el estilo dramático del deck
    divulgativo (slides DEMO: fondo navy 0F2A43, título verde 00E08F 40pt bold,
    cuerpo blanco 26pt).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Fondo navy (igual que las slides DEMO 0F2A43)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x43)

    # Título dramático verde
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12.0), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Le hablas… y obedece"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0x00, 0xE0, 0x8F)

    # Cuerpo blanco 26pt
    body = slide.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(6.0), Inches(3.6))
    btf = body.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = (
        "Pídele «el cubo rojo de la izquierda» y lo hace. El mismo sistema "
        "—ojos, cerebro y brazo— ahora también entiende lenguaje natural."
    )
    br.font.size = Pt(26)
    br.font.bold = False
    br.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Imagen (lado derecho)
    slide.shapes.add_picture(ASSET_IMG, Inches(7.0), Inches(1.9), width=Inches(5.5))

    return slide


# ---------------------------------------------------------------------------
def process_deck(path, marker, builder, target_index, note_text):
    """Procesa un deck: idempotente. Inserta la slide si no existe, refresca
    siempre las notas del ponente del slide-marcador y guarda una sola vez.

    Devuelve dict con info de verificación.
    """
    prs = Presentation(path)
    before_titles = _slide_titles(prs)
    n_before = len(before_titles)

    skipped = _has_marker(prs, marker)
    if skipped:
        print(f"  -> '{marker}' ya existe en {path}, no duplico "
              f"(slides={n_before})")
        n_after = n_before
    else:
        builder(prs)
        _move_last_slide_to(prs, target_index)
        n_after = n_before + 1
        print(f"  -> insertada slide en {path} (posición {target_index}, "
              f"slides {n_before}->{n_after})")

    # Notas del ponente: siempre (sobrescribe, idempotente)
    set_notes(prs, marker, note_text)
    print(f"  -> notas del ponente fijadas en slide-marcador de {path}")

    # Guardado único por deck, ocurra o no la inserción
    prs.save(path)

    return {
        "path": path, "skipped": skipped,
        "n_before": n_before, "n_after": n_after,
        "before_titles": before_titles,
    }


def verify_deck(path, expected_after, target_index, marker_in_title,
                before_titles):
    """Reabre el deck y verifica posición, título, imagen y preservación."""
    prs = Presentation(path)
    titles = _slide_titles(prs)
    n = len(titles)
    assert n == expected_after, f"{path}: esperaba {expected_after} slides, hay {n}"

    new_slide = prs.slides[target_index]
    new_title = ""
    n_pics = 0
    for sh in new_slide.shapes:
        if sh.has_text_frame and not new_title and sh.text_frame.text.strip():
            new_title = sh.text_frame.text.strip()
        if sh.shape_type == 13:  # PICTURE
            n_pics += 1
    assert marker_in_title in new_title, \
        f"{path}: título en idx {target_index} = {new_title!r}"
    assert n_pics == 1, f"{path}: esperaba 1 imagen, hay {n_pics}"

    # Preservación: todos los títulos previos siguen presentes
    new_set = set(titles)
    missing = [t for t in before_titles if t and t not in new_set]
    assert not missing, f"{path}: títulos perdidos: {missing}"

    print(f"  OK {path}: {n} slides; idx {target_index} título={new_title!r}; "
          f"imágenes={n_pics}; preservación OK ({len(before_titles)} títulos)")
    return n, new_title, n_pics


# ---------------------------------------------------------------------------
def _find_marker_slide(prs, marker):
    """Devuelve la primera slide cuyo texto contiene el marcador (o None)."""
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and marker in sh.text_frame.text:
                return s
    return None


def verify_notes(path, marker, expected_count, note_substring):
    """Reabre el deck y verifica recuento de slides + notas del slide-marcador."""
    prs = Presentation(path)
    n = len(list(prs.slides))
    assert n == expected_count, \
        f"{path}: esperaba {expected_count} slides, hay {n}"

    s = _find_marker_slide(prs, marker)
    assert s is not None, f"{path}: no se encontró slide con marcador {marker!r}"
    assert s.has_notes_slide, f"{path}: el slide-marcador no tiene notas"
    notes_text = s.notes_slide.notes_text_frame.text
    assert note_substring in notes_text, \
        f"{path}: notas no contienen {note_substring!r}; notas={notes_text!r}"

    print(f"  OK {path}: {n} slides; notas del slide-marcador contienen "
          f"{note_substring!r}")
    return n, notes_text


def main():
    print("== SP-3: slides de lenguaje natural ==")
    print("[1] Procesando decks")
    process_deck(DECK_DEFENSA, MARKER_DEFENSA, build_defensa_slide,
                 TARGET_INDEX, NOTE_DEFENSA)
    process_deck(DECK_DIVULGA, MARKER_DIVULGA, build_divulga_slide,
                 TARGET_INDEX, NOTE_DIVULGA)

    print("[2] Verificación (reabriendo)")
    verify_notes(DECK_DEFENSA, MARKER_DEFENSA, 19,
                 "selección del objeto correcto acierta el 100 %")
    verify_notes(DECK_DIVULGA, MARKER_DIVULGA, 22, "háblale al robot")

    print("== Hecho ==")


if __name__ == "__main__":
    main()
