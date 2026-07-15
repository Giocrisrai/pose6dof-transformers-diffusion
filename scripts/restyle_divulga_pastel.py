#!/usr/bin/env python3
"""Tema CLARO 'Menta & Coral' para la divulgativa (sobre una COPIA del _base).

Pensado para auditorio de DÍA (el tema oscuro se lava con el proyector):
- Fondo crema, texto carbón legible, títulos coral, acentos menta/coral/butter.
- Barra de acento pastel a la izquierda.
- Tarjetas blancas tras los infográficos (resaltan sobre el crema).
- Botones de demo recoloreados a pastel; hipervínculos y QR intactos.
- MARCA personal en cada slide: monograma 'GG' + 'Giocrisrai Godoy · @giocrisrai'.

NO toca el original: lee SRC (_base) y escribe DST (canónica).
Uso:  ../.venv_thesis/bin/python scripts/restyle_divulga_pastel.py
"""

from __future__ import annotations

import io
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
SRC = REPO / "docs/defensa/Presentacion_Robotica_IA_base.pptx"
DST = REPO / "docs/defensa/Presentacion_Robotica_IA.pptx"

# ---- Paleta Menta & Coral (clara) ----
CREAM = RGBColor(0xFF, 0xFD, 0xF7)  # fondo
INK = RGBColor(0x1F, 0x24, 0x30)  # texto cuerpo (carbón)
INK_SOFT = RGBColor(0x6B, 0x72, 0x82)  # texto secundario
TITLE = RGBColor(0xE8, 0x5D, 0x3D)  # coral profundo (títulos, legible en crema)
TERM = RGBColor(0x0E, 0x8C, 0x73)  # menta profunda (sub-acentos en texto)
MINT = RGBColor(0x7E, 0xD9, 0xC3)  # acento pastel
CORAL = RGBColor(0xFF, 0x9E, 0x80)  # acento pastel
BUTTER = RGBColor(0xFF, 0xD5, 0x6B)  # acento pastel
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# acentos del base que actúan como TÍTULO -> coral
TITLE_ACCENTS = {"00E08F", "34D399", "35876B"}
# acentos del base que actúan como SUB-ACENTO de texto -> menta profunda
TERM_ACCENTS = {"0098CD", "22D3EE", "9FD8EE"}

MARK_BAR = "pastel_bar"
MARK_BRAND = "brand_mark"


def _luma(hexstr: str) -> float:
    try:
        r = int(hexstr[0:2], 16)
        g = int(hexstr[2:4], 16)
        b = int(hexstr[4:6], 16)
        return 0.299 * r + 0.587 * g + 0.114 * b
    except Exception:
        return 0.0


def _run_hex(run):
    try:
        if run.font.color and run.font.color.type is not None:
            return str(run.font.color.rgb).upper()
    except Exception:
        pass
    return None


def set_background(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = CREAM


def restyle_text(slide):
    for sh in slide.shapes:
        if sh.name == MARK_BRAND or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                hexc = _run_hex(run)
                size = run.font.size.pt if run.font.size else 18
                if hexc in TITLE_ACCENTS or size >= 28:
                    run.font.color.rgb = TITLE
                elif hexc in TERM_ACCENTS:
                    run.font.color.rgb = TERM
                else:
                    # cuerpo: carbón (también voltea blanco/claro a oscuro)
                    run.font.color.rgb = INK


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.14), Inches(7.5)
    )
    bar.name = MARK_BAR
    bar.fill.solid()
    bar.fill.fore_color.rgb = MINT
    bar.line.fill.background()
    chip = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0), Inches(0.14), Inches(2.2)
    )
    chip.name = MARK_BAR
    chip.fill.solid()
    chip.fill.fore_color.rgb = CORAL
    chip.line.fill.background()
    spTree = slide.shapes._spTree
    for shp in (bar, chip):
        spTree.remove(shp._element)
        spTree.insert(2, shp._element)


def _is_white_bg(blob) -> bool:
    try:
        from PIL import Image

        im = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = im.size
        pts = [(2, 2), (w - 3, 2), (2, h - 3), (w - 3, h - 3), (w // 2, 2), (w // 2, h - 3)]
        return sum(1 for (x, y) in pts if min(im.getpixel((x, y))) > 235) >= 5  # type: ignore[arg-type,misc]
    except Exception:
        return False


def _soft_shadow(shape):
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(
        qn("a:outerShdw"),
        {"blurRad": "80000", "dist": "30000", "dir": "5400000", "rotWithShape": "0"},
    )
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "9AA0AE"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "45000"}))
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def card_infographics(slide) -> int:
    pics = [
        sh
        for sh in slide.shapes
        if sh.shape_type == 13 and sh.left is not None and _is_white_bg(sh.image.blob)
    ]
    n = 0
    for pic in pics:
        pad = Inches(0.14)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pic.left - pad,
            pic.top - pad,
            pic.width + 2 * pad,
            pic.height + 2 * pad,
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = MINT
        card.line.width = Pt(1.0)
        try:
            _soft_shadow(card)
        except Exception:
            pass
        spTree = slide.shapes._spTree
        spTree.remove(card._element)
        spTree.insert(list(spTree).index(pic._element), card._element)
        n += 1
    return n


def recolor_buttons(slide):
    """Botones de demo (name 'demo_btn') a pastel: relleno menta, texto carbón."""
    for sh in slide.shapes:
        if sh.name == "demo_btn":
            sh.fill.solid()
            sh.fill.fore_color.rgb = MINT
            sh.line.fill.background()
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = INK


def add_brand(slide, page: int):
    """Monograma 'GG' + nombre/handle, abajo-izquierda (marca constante)."""
    mono = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(7.02), Inches(0.42), Inches(0.42)
    )
    mono.name = MARK_BRAND
    mono.fill.solid()
    mono.fill.fore_color.rgb = CORAL
    mono.line.fill.background()
    mp = mono.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "GG"
    mr.font.size = Pt(13)
    mr.font.bold = True
    mr.font.color.rgb = WHITE
    mono.text_frame.margin_top = 0
    mono.text_frame.margin_bottom = 0

    tb = slide.shapes.add_textbox(Inches(0.85), Inches(7.04), Inches(6.5), Inches(0.4))
    tb.name = MARK_BRAND
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Giocrisrai Godoy"
    r1.font.size = Pt(11)
    r1.font.bold = True
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = "   ·   @giocrisrai"
    r2.font.size = Pt(11)
    r2.font.color.rgb = INK_SOFT

    # número de slide, abajo-derecha
    num = slide.shapes.add_textbox(Inches(12.2), Inches(7.04), Inches(0.9), Inches(0.4))
    num.name = MARK_BRAND
    pn = num.text_frame.paragraphs[0]
    pn.alignment = PP_ALIGN.RIGHT
    rn = pn.add_run()
    rn.text = str(page)
    rn.font.size = Pt(11)
    rn.font.color.rgb = INK_SOFT


def add_cover_brand(slide):
    """Marca grande en la portada: monograma + nombre + handles."""
    mono = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.05), Inches(0.7), Inches(0.7)
    )
    mono.name = MARK_BRAND
    mono.fill.solid()
    mono.fill.fore_color.rgb = CORAL
    mono.line.fill.background()
    mp = mono.text_frame.paragraphs[0]
    mp.alignment = PP_ALIGN.CENTER
    mr = mp.add_run()
    mr.text = "GG"
    mr.font.size = Pt(22)
    mr.font.bold = True
    mr.font.color.rgb = WHITE

    tb = slide.shapes.add_textbox(Inches(1.55), Inches(6.02), Inches(9.0), Inches(0.85))
    tb.name = MARK_BRAND
    tf = tb.text_frame
    tf.word_wrap = False
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = "Giocrisrai Godoy"
    r1.font.size = Pt(20)
    r1.font.bold = True
    r1.font.color.rgb = INK
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "linkedin.com/in/giocrisrai   ·   github.com/Giocrisrai"
    r2.font.size = Pt(12)
    r2.font.color.rgb = TERM


def main() -> int:
    shutil.copy(SRC, DST)
    prs = Presentation(DST)
    cards = 0
    for i, slide in enumerate(prs.slides, 1):
        set_background(slide)
        restyle_text(slide)
        add_accent_bar(slide)
        cards += card_infographics(slide)
        recolor_buttons(slide)
        if i == 1:
            add_cover_brand(slide)  # marca grande en portada
        else:
            add_brand(slide, i)  # pie de marca discreto
    prs.save(DST)
    print(f"tarjetas: {cards}")
    print(f"tema 'Menta & Coral' aplicado -> {DST.relative_to(REPO)} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
