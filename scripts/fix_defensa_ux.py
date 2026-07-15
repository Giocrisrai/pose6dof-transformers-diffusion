#!/usr/bin/env python3
"""Correcciones de UX/calidad sobre el _base de la defensa (idempotente).

- P1: el título de la portada (slide 1) coincide EXACTO con el de la tesis.
- P2: el slide 13 (lenguaje natural) baja de 5 a 3 bullets, para que el vídeo
  respire.
- P2: un QR pequeño y discreto al repositorio en el slide de Cierre, en estilo
  UNIR (sin tarjeta; fondo claro).

Tras correr esto, regenera la canónica: scripts/refresh_defensa.py
Uso:  ../.venv_thesis/bin/python scripts/fix_defensa_ux.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parents[1]
DECK = REPO / "docs/defensa/Presentacion_Defensa_TFM_base.pptx"
QR_REPO = REPO / "docs/defensa/assets/qr/qr_github.png"

TITULO_TESIS = ("Estimación de Pose 6-DoF con Transformers y Modelos de Difusión "
                "para Bin Picking Robótico")
GREY = RGBColor(0x55, 0x55, 0x55)

BULLETS_LANG = [
    "«dame el cubo rojo de la izquierda» → selecciona y agarra la pieza descrita",
    "Parser ES/EN + modelo local opcional; anclaje por color/forma/tamaño y posición",
    "Open-license, corre en local; validado E2E (selección 100 %, agarre 4 mm)",
]
QR_NAME = "qr_repo"


def _shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def fix_titulo(prs) -> None:
    s = list(prs.slides)[0]
    tb = _shape(s, "TextBox 1")
    p0 = tb.text_frame.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = TITULO_TESIS
        for r in p0.runs[1:]:
            r.text = ""
    print("  [ok] P1: título de portada alineado con la tesis")


def fix_bullets(prs) -> None:
    s = list(prs.slides)[12]  # slide 13
    tb = _shape(s, "TextBox 3")
    paras = tb.text_frame.paragraphs
    # reescribe los 3 primeros, elimina el resto
    for i, txt in enumerate(BULLETS_LANG):
        p = paras[i]
        if p.runs:
            p.runs[0].text = "•  " + txt
            for r in p.runs[1:]:
                r.text = ""
    for p in list(paras)[len(BULLETS_LANG):]:
        p._p.getparent().remove(p._p)
    print(f"  [ok] P2: slide 13 reducido a {len(BULLETS_LANG)} bullets")


def add_qr_cierre(prs) -> None:
    s = list(prs.slides)[-1]  # Cierre (slide 19)
    if _shape(s, QR_NAME) is not None:
        print("  [skip] QR de repo ya presente en el Cierre")
        return
    pic = s.shapes.add_picture(str(QR_REPO), Inches(11.55), Inches(5.25),
                               width=Inches(1.35), height=Inches(1.35))
    pic.name = QR_NAME
    pic.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    pic.line.width = Pt(0.75)
    cap = s.shapes.add_textbox(Inches(10.9), Inches(6.62), Inches(2.6), Inches(0.35))
    cap.name = QR_NAME
    pcap = cap.text_frame.paragraphs[0]
    pcap.alignment = PP_ALIGN.CENTER
    r = pcap.add_run()
    r.text = "Código y datos abiertos"
    r.font.size = Pt(9)
    r.font.color.rgb = GREY
    print("  [ok] P2: QR del repositorio añadido al Cierre")


def main() -> int:
    prs = Presentation(DECK)
    fix_titulo(prs)
    fix_bullets(prs)
    add_qr_cierre(prs)
    prs.save(DECK)
    print("Ahora regenera la canónica: scripts/refresh_defensa.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
