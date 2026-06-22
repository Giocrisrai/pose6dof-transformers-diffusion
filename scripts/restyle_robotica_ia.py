#!/usr/bin/env python3
"""Restyle 'Vibrante tech' para Presentacion_Robotica_IA.pptx (sobre una COPIA).

Sistema visual:
- Base oscura tech (BG_DARK).
- Acentos vivos: cian eléctrico + lima; magenta puntual.
- Todo el texto a colores claros legibles sobre oscuro; títulos grandes en blanco
  o cian; acentos verdes/azules remapeados a la nueva paleta.
- Barra de acento (gradiente cian→lima) arriba-izquierda en cada slide.

NO toca el original: lee SRC y escribe DST. Las imágenes embebidas se conservan.
Uso:  ../.venv_thesis/bin/python scripts/restyle_robotica_ia.py
"""
from __future__ import annotations

import io
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
SRC = REPO / "docs/entrega3/Presentacion_Robotica_IA.pptx"
DST = REPO / "docs/entrega3/Presentacion_Robotica_IA_v2.pptx"

# ---- Paleta Vibrante tech ----
BG_DARK = RGBColor(0x0B, 0x12, 0x26)     # navy casi negro
CYAN = RGBColor(0x22, 0xD3, 0xEE)        # cian eléctrico
LIME = RGBColor(0x34, 0xD3, 0x99)        # verde lima
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY = RGBColor(0xE2, 0xE8, 0xF0)        # gris claro para cuerpo
MUTED = RGBColor(0x94, 0xA3, 0xB8)       # gris medio para secundario

# colores "oscuros" del deck original que hay que aclarar sobre fondo oscuro
DARKISH = {"0F2A43", "333333", "666666", "222222", "2C3E50", "000000", "1F4E79"}
# remapeo de acentos del deck a la nueva paleta
REMAP = {"9FD8EE": CYAN, "0098CD": CYAN, "00E08F": LIME, "35876B": LIME}


def _luma_is_dark(hexstr: str) -> bool:
    try:
        r = int(hexstr[0:2], 16); g = int(hexstr[2:4], 16); b = int(hexstr[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) < 140
    except Exception:
        return False


def _run_hex(run):
    try:
        if run.font.color and run.font.color.type is not None:
            return str(run.font.color.rgb)
    except Exception:
        pass
    return None


def set_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_accent_bar(slide):
    """Barra de acento corta arriba-izquierda (cian) + tope lima, look dinámico."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0),
                                 Inches(0.12), Inches(7.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN
    bar.line.fill.background()
    chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.0), Inches(0.0),
                                  Inches(0.12), Inches(2.2))
    chip.fill.solid(); chip.fill.fore_color.rgb = LIME
    chip.line.fill.background()
    # enviar al fondo para no tapar contenido
    spTree = slide.shapes._spTree
    for shp in (bar, chip):
        spTree.remove(shp._element)
        spTree.insert(2, shp._element)


def restyle_text(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                hexc = _run_hex(run)
                size = run.font.size.pt if run.font.size else 18
                if hexc in REMAP:
                    run.font.color.rgb = REMAP[hexc]
                elif hexc == "FFFFFF":
                    run.font.color.rgb = WHITE                      # ya claro
                elif hexc is None or hexc.upper() in DARKISH or _luma_is_dark(hexc or ""):
                    # texto oscuro/temático -> claro; títulos grandes en blanco
                    run.font.color.rgb = WHITE if size >= 28 else BODY
                else:
                    # color claro desconocido -> dejar, pero asegurar legibilidad
                    if _luma_is_dark(hexc):
                        run.font.color.rgb = BODY


def _is_white_bg(blob) -> bool:
    """True si la imagen tiene fondo claro (infográfico) según muestreo de bordes."""
    try:
        from PIL import Image
        im = Image.open(io.BytesIO(blob)).convert("RGB")
        w, h = im.size
        pts = [(2, 2), (w - 3, 2), (2, h - 3), (w - 3, h - 3), (w // 2, 2), (w // 2, h - 3)]
        return sum(1 for (x, y) in pts if min(im.getpixel((x, y))) > 235) >= 5
    except Exception:
        return False


def _add_soft_shadow(shape) -> None:
    """Sombra exterior suave vía XML (a:outerShdw)."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(qn("a:outerShdw"),
                           {"blurRad": "90000", "dist": "38000",
                            "dir": "5400000", "rotWithShape": "0"})
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "40000"}))
    shdw.append(clr); eff.append(shdw); spPr.append(eff)


def card_infographics(slide) -> int:
    """Pone una tarjeta clara redondeada (con sombra) detrás de cada infográfico
    de fondo blanco, para que se lea como figura intencional sobre el deck oscuro.
    No toca fotos del sim ni el reel. Devuelve nº de tarjetas añadidas."""
    pics = [sh for sh in slide.shapes
            if sh.shape_type == 13 and sh.left is not None
            and _is_white_bg(sh.image.blob)]
    n = 0
    for pic in pics:
        pad = Inches(0.14)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            pic.left - pad, pic.top - pad, pic.width + 2 * pad, pic.height + 2 * pad)
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFB)
        card.line.color.rgb = RGBColor(0x1E, 0x29, 0x3B); card.line.width = Pt(0.75)
        try:
            _add_soft_shadow(card)
        except Exception:
            pass
        # mover la tarjeta justo DETRÁS de la imagen (mismo z, una posición antes)
        spTree = slide.shapes._spTree
        spTree.remove(card._element)
        spTree.insert(list(spTree).index(pic._element), card._element)
        n += 1
    return n


def main() -> int:
    shutil.copy(SRC, DST)
    prs = Presentation(DST)
    total_cards = 0
    for slide in prs.slides:
        set_background(slide)
        restyle_text(slide)
        add_accent_bar(slide)
        total_cards += card_infographics(slide)
    prs.save(DST)
    print(f"tarjetas de figura añadidas: {total_cards}")
    print(f"restyle aplicado -> {DST.relative_to(REPO)} ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
