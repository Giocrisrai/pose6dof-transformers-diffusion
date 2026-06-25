#!/usr/bin/env python3
"""Añade a la divulgativa (_base) dos slides: glosario 'Para entendernos' y
'Conéctate y úsalo' (QR de LinkedIn, PyPI y GitHub + email + pip install).

Diseño:
- Opera sobre Presentacion_Robotica_IA_base.pptx; la canónica se regenera con
  restyle_divulga_pastel.py (tema claro 'Menta & Coral').
- Colores elegidos para que el restyle claro los mapee bien:
    título 00E08F -> coral, término 0098CD -> menta, cuerpo 222222 -> carbón.
- Los QR (PNG blanco) reciben automáticamente la tarjeta blanca del restyle.
- Idempotente: si ya existe un slide con el marcador (en el título), no recrea.
- DEBE correr DESPUÉS de embed_sim_videos.py (ese usa índices fijos sobre el
  base pristino); este inserta slides y desplaza los índices.

Uso:  ../.venv_thesis/bin/python scripts/build_divulga_extras.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO = Path(__file__).resolve().parents[1]
DECK = REPO / "docs/entrega3/Presentacion_Robotica_IA_base.pptx"
QR = REPO / "docs/entrega3/assets/qr"

TITLE_GREEN = RGBColor(0x00, 0xE0, 0x8F)   # -> LIME tras restyle
TERM_BLUE = RGBColor(0x00, 0x98, 0xCD)     # -> CYAN tras restyle
BODY = RGBColor(0x22, 0x22, 0x22)          # -> claro tras restyle
NAVY = RGBColor(0x0F, 0x2A, 0x43)          # bg base (restyle lo oscurece)

MARK_GLOSARIO = "Para entendernos"
MARK_CONTACTO = "Conéctate y úsalo"

CONCEPTOS = [
    ("Pose 6-DoF",
     "dónde está algo y cómo está girado: 3 de posición + 3 de giro."),
    ("Transformer",
     "la familia de ChatGPT; aprende a prestar atención a lo importante."),
    ("Modelo de difusión",
     "la familia de los generadores de imágenes: del ruido a la respuesta."),
    ("Visual servoing",
     "el robot se corrige mirando, como tu mano al enhebrar una aguja."),
    ("Bin-picking",
     "sacar piezas amontonadas de una caja."),
]


def _titles(prs):
    out = []
    for s in prs.slides:
        t = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().split("\n")[0]
                break
        out.append(t)
    return out


def _has(prs, mark):
    return any(mark in t for t in _titles(prs))


def _bg_navy(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = NAVY


def _title(slide, text, size=40):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = TITLE_GREEN


def _move_last_to(prs, idx):
    lst = prs.slides._sldIdLst
    new = list(lst)[-1]
    lst.remove(new)
    lst.insert(idx, new)


def build_glosario(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_navy(slide)
    _title(slide, MARK_GLOSARIO)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.85), Inches(11.8), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (term, expl) in enumerate(CONCEPTOS):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        rt = p.add_run()
        rt.text = f"{term}:  "
        rt.font.size = Pt(22)
        rt.font.bold = True
        rt.font.color.rgb = TERM_BLUE
        re = p.add_run()
        re.text = expl
        re.font.size = Pt(20)
        re.font.color.rgb = BODY
    return slide


def build_contacto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg_navy(slide)
    _title(slide, MARK_CONTACTO)

    # CTA
    cta = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(12.0), Inches(0.7))
    r = cta.text_frame.paragraphs[0].add_run()
    r.text = "Pruébalo, dame feedback y mejorémoslo juntos."
    r.font.size = Pt(22)
    r.font.color.rgb = BODY

    # 3 QR en fila (clicables: la imagen y el caption enlazan a su URL)
    items = [
        (QR / "qr_linkedin.png", "LinkedIn", "Hablemos",
         "https://www.linkedin.com/in/giocrisrai/"),
        (QR / "qr_pypi.png", "pip install bop-bootstrap-ci", "úsalo en tus proyectos",
         "https://pypi.org/project/bop-bootstrap-ci/"),
        (QR / "qr_github.png", "GitHub", "código completo del pipeline",
         "https://github.com/Giocrisrai/pose6dof-transformers-diffusion"),
    ]
    qr_w = Inches(2.4)
    xs = [Inches(1.35), Inches(5.45), Inches(9.55)]
    top = Inches(2.55)
    from pptx.enum.text import PP_ALIGN
    for (img, cap1, cap2, url), x in zip(items, xs):
        pic = slide.shapes.add_picture(str(img), x, top, width=qr_w, height=qr_w)
        pic.click_action.hyperlink.address = url      # QR clicable
        # caption (centrada bajo el QR)
        cb = slide.shapes.add_textbox(x - Inches(0.5), top + qr_w + Inches(0.08),
                                      qr_w + Inches(1.0), Inches(1.1))
        ctf = cb.text_frame
        ctf.word_wrap = True
        p1 = ctf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = cap1
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = TERM_BLUE
        r1.hyperlink.address = url                     # caption clicable
        p2 = ctf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = cap2
        r2.font.size = Pt(13)
        r2.font.color.rgb = BODY

    # email al pie
    em = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.5))
    from pptx.enum.text import PP_ALIGN
    pe = em.text_frame.paragraphs[0]
    pe.alignment = PP_ALIGN.CENTER
    re = pe.add_run()
    re.text = "Giocrisrai Godoy · giocrisrai@gmail.com"
    re.font.size = Pt(15)
    re.font.color.rgb = BODY
    return slide


def main() -> int:
    """Inserta los dos slides si no existen. Para editar su contenido, parte de
    un _base SIN extras (p.ej. restaurándolo de git) y vuelve a ejecutar; así se
    evitan partes huérfanas dentro del .pptx."""
    prs = Presentation(DECK)
    n0 = len(prs.slides._sldIdLst)

    if not _has(prs, MARK_GLOSARIO):
        build_glosario(prs)
        _move_last_to(prs, 9)   # tras el slide 9 -> nuevo slide 10
        print("  [ok] insertado 'Para entendernos' como slide 10")
    else:
        print("  [skip] 'Para entendernos' ya existe (restaura el _base para editar)")

    if not _has(prs, MARK_CONTACTO):
        build_contacto(prs)     # queda al final
        print("  [ok] añadido 'Conéctate y úsalo' al final")
    else:
        print("  [skip] 'Conéctate y úsalo' ya existe (restaura el _base para editar)")

    prs.save(DECK)
    n1 = len(prs.slides._sldIdLst)
    print(f"  slides: {n0} -> {n1}")
    print("Ahora regenera la canónica: scripts/restyle_divulga_pastel.py")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
